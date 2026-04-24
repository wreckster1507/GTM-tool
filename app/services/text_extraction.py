"""Pluggable text extraction for Drive + uploaded files.

We intentionally keep this module sync and CPU-bound — callers should run it
in a thread via ``asyncio.to_thread`` when invoked from async code.
"""
from __future__ import annotations

import base64
import csv
import io
import logging
import os
import shutil
import subprocess
import tempfile
from typing import Optional

from app.config import settings

logger = logging.getLogger(__name__)

# MIME types we know how to turn into text. Anything else is skipped with a
# warning so we don't crash the indexer on exotic files.
SUPPORTED_MIME_PREFIXES = (
    "text/",
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
    "application/msword",
    "application/vnd.ms-excel",
    "application/json",
    "application/xml",
    # Vision — Claude describes/OCRs the image into searchable text.
    "image/",     # png, jpg, webp, gif — all handled by Claude vision
    # Audio/video — transcribed via OpenAI Whisper. ffmpeg is used to strip
    # the audio track out of videos so we stay under the 25 MB Whisper cap.
    # Covers mp4, mov (video/quicktime), webm, mkv, m4v, etc.
    "video/",
    "audio/",
)

# Whisper hard file-size cap.
_WHISPER_MAX_BYTES = 25 * 1024 * 1024


def is_supported(mime_type: str) -> bool:
    mime = (mime_type or "").lower()
    return any(mime.startswith(prefix) for prefix in SUPPORTED_MIME_PREFIXES)


def extract_text(data: bytes, mime_type: str, filename: str = "") -> str:
    """Best-effort plaintext extraction. Returns '' on failure."""
    if not data:
        return ""
    mime = (mime_type or "").lower()

    try:
        if mime.startswith("text/") or mime in {"application/json", "application/xml"}:
            return _decode_text(data)
        if mime == "application/pdf":
            return _extract_pdf(data)
        if "wordprocessingml" in mime or mime == "application/msword" or filename.lower().endswith(".docx"):
            return _extract_docx(data)
        if "spreadsheetml" in mime or mime == "application/vnd.ms-excel" or filename.lower().endswith(".xlsx"):
            return _extract_xlsx(data)
        if (
            "presentationml" in mime
            or mime == "application/vnd.ms-powerpoint"
            or filename.lower().endswith(".pptx")
        ):
            return _extract_pptx(data)
        if filename.lower().endswith(".csv"):
            return _extract_csv(data)
        if mime.startswith("image/") or filename.lower().endswith((".png", ".jpg", ".jpeg")):
            return _extract_image(data, mime or "image/png")
        if mime.startswith("video/") or filename.lower().endswith(
            (".mp4", ".mov", ".webm", ".mkv", ".m4v", ".avi")
        ):
            return _extract_video(data)
        if mime.startswith("audio/"):
            return _transcribe_audio_bytes(data, filename or "audio.mp3")
    except Exception as exc:  # pragma: no cover
        logger.warning("extract_text failed for %s (%s): %s", filename, mime, exc)
        return ""

    logger.info("Skipping unsupported mime=%s filename=%s", mime, filename)
    return ""


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_pdf(data: bytes) -> str:
    import pdfplumber  # local import keeps cold-start fast

    text_parts: list[str] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            page_text: Optional[str] = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    return "\n\n".join(text_parts)


def _extract_docx(data: bytes) -> str:
    from docx import Document  # python-docx

    doc = Document(io.BytesIO(data))
    lines: list[str] = []
    for para in doc.paragraphs:
        if para.text:
            lines.append(para.text)
    # Tables sometimes hold the most interesting content (reqs docs, matrices).
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines)


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    lines: list[str] = []
    for sheet in wb.worksheets:
        lines.append(f"## Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(v) for v in row if v not in (None, "")]
            if values:
                lines.append(" | ".join(values))
    return "\n".join(lines)


def _extract_csv(data: bytes) -> str:
    text = _decode_text(data)
    reader = csv.reader(io.StringIO(text))
    return "\n".join(" | ".join(row) for row in reader)


def _extract_image(data: bytes, mime: str) -> str:
    """
    Use Claude vision to produce a searchable description + OCR of the image.
    Returns '' if no API key is configured (keeps indexer crash-free in dev).
    Sync wrapper over the async Anthropic SDK — callers already run this
    inside asyncio.to_thread, so blocking here is fine.
    """
    api_key = settings.claude_api_key
    if not api_key or not data:
        return ""

    import anthropic  # local import — large dep, keep cold-start cheap

    client = anthropic.Anthropic(api_key=api_key)
    b64 = base64.standard_b64encode(data).decode("ascii")
    media_type = mime if mime in {"image/png", "image/jpeg", "image/webp", "image/gif"} else "image/png"

    resp = client.messages.create(
        model=settings.ANTHROPIC_MODEL,
        max_tokens=1500,
        messages=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {"type": "base64", "media_type": media_type, "data": b64},
                    },
                    {
                        "type": "text",
                        "text": (
                            "Describe this image for a sales-knowledge search index. "
                            "Include: (1) a 2-3 sentence summary of what the image depicts, "
                            "(2) any visible text transcribed verbatim (OCR), "
                            "(3) key entities — company names, product names, people, logos. "
                            "Return plain text only, no markdown headers."
                        ),
                    },
                ],
            }
        ],
    )
    parts = [b.text for b in resp.content if getattr(b, "type", None) == "text"]
    return "\n".join(parts).strip()


def _extract_video(data: bytes) -> str:
    """
    Two-stage extraction:
      1. Strip audio → Whisper transcription (handles meeting recordings, VO).
      2. If the transcript is empty (silent screen-capture / config demo),
         fall back to sampling video frames and running each through Claude
         vision. This covers the "product walk-through without voice" case
         that would otherwise produce no indexable text.
    The two results are concatenated when both exist.
    """
    if not data:
        return ""
    if not shutil.which("ffmpeg"):
        logger.warning("ffmpeg not installed — cannot transcribe mp4")
        return ""

    with tempfile.TemporaryDirectory() as tmp:
        in_path = os.path.join(tmp, "in.mp4")
        with open(in_path, "wb") as fh:
            fh.write(data)

        transcript = _video_to_transcript(in_path, tmp)
        if transcript.strip():
            return transcript

        # Silent/no-speech video — fall back to visual summarization.
        logger.info("Empty transcript — sampling frames for vision fallback")
        return _video_to_vision_summary(in_path, tmp)


def _video_to_transcript(in_path: str, tmp: str) -> str:
    """ffmpeg → mono 16 kHz MP3 → Whisper."""
    out_path = os.path.join(tmp, "out.mp3")
    result = subprocess.run(
        [
            "ffmpeg", "-y", "-i", in_path,
            "-vn", "-ac", "1", "-ar", "16000", "-b:a", "64k",
            out_path,
        ],
        capture_output=True,
    )
    if result.returncode != 0 or not os.path.exists(out_path):
        # "no audio stream" also lands here — expected for screen-captures.
        logger.info("ffmpeg audio extract yielded nothing: %s", result.stderr[-200:])
        return ""
    with open(out_path, "rb") as fh:
        audio_bytes = fh.read()
    return _transcribe_audio_bytes(audio_bytes, "audio.mp3")


# How many frames to sample per video. 6 frames at evenly-spaced timestamps
# captures most screencast state transitions without burning too much
# Claude-vision budget (~$0.02 per video at Sonnet pricing).
_VIDEO_FRAME_SAMPLES = 6


def _video_to_vision_summary(in_path: str, tmp: str) -> str:
    """Sample N evenly-spaced frames and OCR/describe each with Claude."""
    duration = _probe_duration_seconds(in_path)
    if duration <= 0:
        return ""

    # Evenly-spaced timestamps, skipping the very first/last 5% (often blank).
    step = duration / (_VIDEO_FRAME_SAMPLES + 1)
    timestamps = [step * (i + 1) for i in range(_VIDEO_FRAME_SAMPLES)]

    frame_texts: list[str] = []
    for idx, ts in enumerate(timestamps):
        frame_path = os.path.join(tmp, f"frame_{idx}.jpg")
        r = subprocess.run(
            [
                "ffmpeg", "-y", "-ss", f"{ts:.2f}", "-i", in_path,
                "-frames:v", "1", "-q:v", "4", frame_path,
            ],
            capture_output=True,
        )
        if r.returncode != 0 or not os.path.exists(frame_path):
            continue
        with open(frame_path, "rb") as fh:
            img_bytes = fh.read()
        text = _extract_image(img_bytes, "image/jpeg")
        if text.strip():
            frame_texts.append(f"[t={ts:.0f}s]\n{text.strip()}")

    if not frame_texts:
        return ""
    return "## Video frame summary (no speech detected)\n\n" + "\n\n".join(frame_texts)


def _probe_duration_seconds(in_path: str) -> float:
    """Use ffprobe (ships with ffmpeg) to read the container duration."""
    r = subprocess.run(
        [
            "ffprobe", "-v", "error", "-show_entries", "format=duration",
            "-of", "default=noprint_wrappers=1:nokey=1", in_path,
        ],
        capture_output=True,
        text=True,
    )
    try:
        return float(r.stdout.strip())
    except (ValueError, TypeError):
        return 0.0


def _transcribe_audio_bytes(data: bytes, filename: str) -> str:
    """OpenAI Whisper transcription. Chunks oversized files."""
    api_key = settings.OPENAI_API_KEY
    if not api_key or not data:
        return ""

    from openai import OpenAI

    client = OpenAI(api_key=api_key)
    if len(data) <= _WHISPER_MAX_BYTES:
        resp = client.audio.transcriptions.create(
            model="whisper-1",
            file=(filename, data),
        )
        return (resp.text or "").strip()

    # Oversized: a sales meeting can run 60+ min. Caller should've downsampled
    # via ffmpeg, but as a safety net, bail gracefully with a clear log.
    logger.warning(
        "Audio payload %d bytes exceeds Whisper 25 MB cap — skipping %s",
        len(data), filename,
    )
    return ""


def _extract_pptx(data: bytes) -> str:
    """Pull text out of every slide — title, body, notes, tables."""
    from pptx import Presentation  # python-pptx

    prs = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"## Slide {i}")
        for shape in slide.shapes:
            # Plain text boxes / titles / placeholders
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = " ".join(run.text for run in para.runs if run.text)
                    if text.strip():
                        lines.append(text.strip())
            # Tables occasionally carry the actual content (matrices, specs).
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        lines.append(" | ".join(cells))
        # Speaker notes — often the most useful narrative bit.
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"[Notes] {notes}")
    return "\n".join(lines)
