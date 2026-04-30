"""PoC Demo PPT generator — fills Beacon's Drive PoC Demo deck template.

Mirrors the PoC Kickoff generator in shape: discover the indexed Drive
template by name keywords, export the Google Slides deck to .pptx, hand
the slide-by-slide structure to Claude with the ``poc_demo_ppt``
instruction set, patch rewrites back, upload as a new editable Google
Slides deck.

Slide contract (per the Zellis template):
  Slides 1, 2, 6, 7  — STATIC (Beacon platform pitch / demo holder /
                       roadmap). Claude returns these unchanged.
  Slides 3, 4, 5     — FILLED from email + PoC Kickoff content. Slide 5's
                       heading still reads "Zellis" in the template; the
                       rewriter replaces it with the actual client name.
"""
from __future__ import annotations

import io
import logging
from dataclasses import dataclass
from datetime import datetime
from typing import Optional

from sqlmodel import select as sm_select

from app.clients.anthropic_client import get_anthropic_client
from app.clients.google_drive import upload_as_google_slides
from app.config import settings
from app.database import AsyncSessionLocal as async_session
from app.models.user_email_connection import UserEmailConnection
from app.models.zippy import IndexedDriveFile
from app.services.zippy_docs.base import (
    GeneratedDocument,
    build_output_path,
    human_today,
)
from app.services.zippy_docs.claude_rewriter import rewrite_with_claude
from app.services.zippy_docs.doc_rewriter import (
    extract_pptx_structure,
    rewrite_pptx_content,
)

logger = logging.getLogger(__name__)


@dataclass
class PoCPPTInput:
    """Inputs for the PoC Demo PPT generator.

    Two content sources are supported because PoC Demo decks usually
    summarise both the kickoff (scope, objective) and post-kickoff
    progress emails (capabilities demonstrated, impact). Either may be
    empty; Claude will fill what it can and leave the rest as the
    template's Zellis examples (with the company name swapped in).
    """
    client_name: str
    email_thread_content: str = ""
    poc_kickoff_content: str = ""
    meeting_date: Optional[str] = None
    extra_context: Optional[str] = None


# ── Drive template discovery ────────────────────────────────────────────────

_POC_PPT_KEYWORDS = [
    "zellis - poc ppt",
    "poc ppt",
    "poc demo ppt",
    "poc presentation",
    "poc demo presentation",
    "poc ppt template",
]

# Per-file in-memory cache, keyed by drive_modified_at — matches the
# pattern used by mom.py / poc_kickoff.py so a re-uploaded template
# invalidates automatically.
_TEMPLATE_CACHE: dict[str, tuple[str, bytes]] = {}


async def _find_poc_ppt_template_row(
    user_id: Optional[str] = None,
) -> Optional[IndexedDriveFile]:
    """Return the IndexedDriveFile row for the PoC PPT template, or None.

    Owner priority mirrors poc_kickoff.py: prefer the user's own synced
    file, fall back to admin/shared, never leak another user's private file.
    """
    from sqlalchemy import or_, case as sa_case

    async with async_session() as session:
        for keyword in _POC_PPT_KEYWORDS:
            stmt = sm_select(IndexedDriveFile).where(
                IndexedDriveFile.name.ilike(f"%{keyword}%"),
            )
            if user_id is not None:
                stmt = stmt.where(
                    or_(
                        IndexedDriveFile.owner_user_id == user_id,
                        IndexedDriveFile.is_admin == True,  # noqa: E712
                    )
                )
            else:
                stmt = stmt.where(IndexedDriveFile.is_admin == True)  # noqa: E712

            if user_id is not None:
                user_priority = sa_case(
                    (IndexedDriveFile.owner_user_id == user_id, 0),
                    else_=1,
                )
                stmt = stmt.order_by(
                    user_priority,
                    IndexedDriveFile.last_indexed_at.desc(),
                )
            else:
                stmt = stmt.order_by(IndexedDriveFile.last_indexed_at.desc())

            row = (await session.execute(stmt.limit(1))).scalar_one_or_none()
            if row:
                logger.info(
                    "PoC PPT template found via keyword=%r: name=%s "
                    "file_id=%s mime=%s owner=%s is_admin=%s",
                    keyword, row.name, row.drive_file_id, row.mime_type,
                    row.owner_user_id, row.is_admin,
                )
                return row
    logger.info("No PoC PPT template matched for user_id=%s", user_id)
    return None


async def _fetch_template_bytes(
    user_id: Optional[str] = None,
) -> tuple[Optional[bytes], Optional[str]]:
    """Fetch the PoC PPT template bytes for the given user.

    Google Slides → export as .pptx. Falls back to a direct download
    when the template is uploaded as a real .pptx file.
    """
    row = await _find_poc_ppt_template_row(user_id=user_id)
    if not row:
        return None, (
            "No PoC PPT template found. Add a file with 'POC PPT' in the "
            "name to your Drive folder and re-sync."
        )

    file_id = row.drive_file_id
    actual_mime = row.mime_type
    cache_key = (
        row.drive_modified_at.isoformat()
        if row.drive_modified_at else "no-modified"
    )

    cached = _TEMPLATE_CACHE.get(file_id)
    if cached and cached[0] == cache_key:
        logger.info(
            "PoC PPT template cache HIT for %s (%d bytes)",
            file_id, len(cached[1]),
        )
        return cached[1], None

    async with async_session() as session:
        result = await session.execute(
            sm_select(UserEmailConnection).where(
                UserEmailConnection.user_id == row.owner_user_id,
                UserEmailConnection.is_active == True,  # noqa: E712
            )
        )
        connection = result.scalar_one_or_none()

    if not connection:
        return None, (
            f"No active Drive connection for the owner of '{row.name}'. "
            "Reconnect Google Drive in Settings and re-sync."
        )

    try:
        if actual_mime == "application/vnd.google-apps.presentation":
            # Google Slides → export to .pptx so doc_rewriter can walk
            # the deck with python-pptx. Same export-endpoint shape as
            # the Google-Doc branch in poc_kickoff.py.
            import httpx
            from app.clients.google_drive import (
                _ensure_token,
                DRIVE_API_BASE,
            )
            access_token, _updated = await _ensure_token(
                connection.token_data,
                settings.gmail_client_id,
                settings.gmail_client_secret,
            )
            pptx_mime = (
                "application/vnd.openxmlformats-officedocument."
                "presentationml.presentation"
            )
            async with httpx.AsyncClient(timeout=60) as http:
                resp = await http.get(
                    f"{DRIVE_API_BASE}/files/{file_id}/export",
                    headers={"Authorization": f"Bearer {access_token}"},
                    params={"mimeType": pptx_mime},
                )
            if resp.status_code != 200:
                return None, (
                    f"Drive export to .pptx failed for '{row.name}' "
                    f"(status {resp.status_code}): {resp.text[:300]}"
                )
            data = resp.content
        else:
            # Already a .pptx in Drive — straight download.
            from app.clients import google_drive
            data, _mime, _updated = await google_drive.download_file_bytes(
                file_id=file_id,
                mime_type=actual_mime,
                token_data=connection.token_data,
                client_id=settings.gmail_client_id,
                client_secret=settings.gmail_client_secret,
            )
        if not data:
            return None, (
                f"Drive returned empty bytes for {row.name} "
                f"(file_id={file_id}, mime={actual_mime})."
            )
        if not data.startswith(b"PK"):
            return None, (
                f"'{row.name}' didn't come back as a valid .pptx "
                f"(first bytes: {data[:8]!r})."
            )
        logger.info("PoC PPT template downloaded fresh: %d bytes", len(data))
        _TEMPLATE_CACHE[file_id] = (cache_key, data)
        return data, None
    except Exception as exc:
        logger.exception(
            "Failed to download PoC PPT template file_id=%s: %s",
            file_id, exc,
        )
        return None, f"Download failed: {exc}"


# ── Public entry points ─────────────────────────────────────────────────────


async def inspect_poc_ppt_template(
    user_id: Optional[str] = None,
) -> dict:
    """Confirm the PoC PPT template is reachable and report slide layout."""
    template_bytes, err = await _fetch_template_bytes(user_id=user_id)
    if not template_bytes:
        return {"found": False, "error": err}

    structure = extract_pptx_structure(template_bytes)
    slide_count = max(
        (b.get("slide_index", 0) for b in structure),
        default=0,
    ) + 1
    row = await _find_poc_ppt_template_row(user_id=user_id)
    return {
        "found": True,
        "template_name": row.name if row else "POC PPT Template",
        "slide_count": slide_count,
        "fillable_slides": [3, 4, 5],
        "note": (
            "Slides 3-5 are filled from email/PoC Kickoff content. "
            "Slides 1, 2, 6, 7 are static."
        ),
    }


async def generate(
    data: PoCPPTInput,
    user_id: Optional[str] = None,
) -> GeneratedDocument:
    """Produce the filled PoC Demo .pptx (and upload as Google Slides)."""
    path, url = build_output_path(
        "PoC_PPT", data.client_name, extension="pptx"
    )
    template_bytes, err = await _fetch_template_bytes(user_id=user_id)
    source = "template"

    if not template_bytes:
        logger.warning("PPT template not found: %s", err)
        _render_fallback_pptx(data, path)
        source = "fallback"
    else:
        structure = extract_pptx_structure(template_bytes)

        # Combine both content sources into one blob — Claude is best
        # given one stream of text rather than juggling two parameters.
        # The headers help it attribute facts ("the kickoff said X,
        # the email said Y") when they conflict.
        combined_content = ""
        if data.poc_kickoff_content:
            combined_content += (
                "=== POC KICKOFF DOCUMENT ===\n"
                f"{data.poc_kickoff_content}\n\n"
            )
        if data.email_thread_content:
            combined_content += (
                "=== EMAIL THREADS ===\n"
                f"{data.email_thread_content}\n\n"
            )

        user_inputs = {
            "client_name": data.client_name,
            "date": data.meeting_date or human_today(),
            "content": combined_content[:25000],
            "extra_context": data.extra_context or "",
        }

        client = get_anthropic_client()
        rewritten = await rewrite_with_claude(
            structure=structure,
            user_inputs=user_inputs,
            doc_type="poc_demo_ppt",
            client=client,
            model=settings.CLAUDE_MODEL_STANDARD,
        )

        out_bytes = rewrite_pptx_content(template_bytes, rewritten)
        path.write_bytes(out_bytes)

    result = GeneratedDocument(
        filename=path.name,
        path=str(path),
        url=url,
        kind="poc_ppt",
        summary=f"PoC Demo PPT for {data.client_name} — {source}.",
        created_at=datetime.utcnow(),
    )

    await _try_upload_to_slides(
        result, path, data.client_name, user_id=user_id
    )
    return result


# ── Fallback ────────────────────────────────────────────────────────────────


def _render_fallback_pptx(data: PoCPPTInput, path) -> None:
    """Minimal placeholder deck when the template can't be fetched.

    Five title-and-content slides — enough for the AE to see Zippy
    didn't silently fail and to know which sections need replacement
    once the real template is restored.
    """
    from pptx import Presentation

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    for title_text, body_text in [
        (
            "AI Implementation Orchestration",
            f"PoC Demo — {data.client_name}",
        ),
        (
            "POC Overview",
            "Content from email threads not yet filled.",
        ),
        (
            "Scope of POC",
            "Please add a POC PPT template to Drive and re-sync.",
        ),
        (
            f"What We've Done on {data.client_name}",
            "",
        ),
        (
            "Proposed Next Steps",
            "PoC Validation → ROI Workshop → Go Live",
        ),
    ]:
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = title_text
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = body_text

    prs.save(str(path))


# ── Drive upload ────────────────────────────────────────────────────────────


async def _try_upload_to_slides(
    doc: GeneratedDocument,
    path,
    client_name: str,
    *,
    user_id: Optional[str],
) -> None:
    """Upload the filled .pptx to Drive as an editable Google Slides deck.

    Best-effort, silently skipped on permission/upload failure — same
    contract as the other generators' upload helpers. Honours the same
    same-day dedup cache so a re-run reuses the existing Slides URL.
    """
    if doc.drive_url:
        return
    if user_id:
        from app.services.zippy_docs.base import (
            cache_upload,
            get_cached_upload,
        )
        cached = get_cached_upload(str(user_id), client_name, doc.kind)
        if cached:
            doc.drive_url = cached
            logger.info(
                "Reusing cached PoC PPT Slides upload: %s", cached
            )
            return
    try:
        from sqlalchemy import or_, case as sa_case

        async with async_session() as session:
            stmt = sm_select(UserEmailConnection).where(
                UserEmailConnection.is_active == True,  # noqa: E712
            )
            if user_id:
                stmt = stmt.where(
                    or_(
                        UserEmailConnection.user_id == user_id,
                        UserEmailConnection.is_admin_folder == True,  # noqa: E712
                    )
                )
                priority = sa_case(
                    (UserEmailConnection.user_id == user_id, 0), else_=1
                )
                stmt = stmt.order_by(priority)
            result = await session.execute(stmt.limit(1))
            connection = result.scalar_one_or_none()

        if not connection:
            logger.info(
                "No active Drive connection — skipping PoC PPT upload"
            )
            return

        pptx_bytes = path.read_bytes()
        slides_name = (
            f"PoC PPT — {client_name} — "
            f"{doc.created_at.strftime('%d %b %Y')}"
        )
        folder_id = getattr(connection, "selected_drive_folder_id", None)

        file_id, web_view_link = await upload_as_google_slides(
            filename=slides_name,
            docx_bytes=pptx_bytes,
            token_data=connection.token_data,
            client_id=settings.gmail_client_id,
            client_secret=settings.gmail_client_secret,
            parent_folder_id=folder_id,
        )
        doc.drive_file_id = file_id
        doc.drive_url = web_view_link
        logger.info("PPT uploaded to Google Slides: %s", web_view_link)
        if user_id and web_view_link:
            from app.services.zippy_docs.base import cache_upload
            cache_upload(
                str(user_id), client_name, doc.kind, web_view_link
            )
    except PermissionError as exc:
        logger.info(
            "drive.file scope not yet granted — skipping PoC PPT "
            "upload: %s", exc,
        )
    except Exception as exc:
        logger.warning(
            "Google Slides upload failed for PoC PPT (non-fatal): %s",
            exc,
        )
