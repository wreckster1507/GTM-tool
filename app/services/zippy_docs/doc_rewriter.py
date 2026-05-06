"""Structure-preserving rewriters for .docx and .pptx templates.

The previous generation relied on ``{{TOKEN}}`` placeholders hand-added to
every template — fragile, and silently no-op when the template file had no
tokens. This module replaces that approach with full content rewriting:

    1. Extract every text-bearing block in document order, flagging which
       ones are structural (headings, short labels) vs. content.
    2. Hand the list to an LLM that rewrites content blocks using the
       user's inputs.
    3. Walk the original document in the same order and overwrite each
       paragraph's text, preserving run-level formatting.

Nothing here knows about MOM/NDA/etc. — the generators call in with the
template bytes, get a structure back, hand the structure to Claude, and pass
the rewritten blocks back through.
"""
from __future__ import annotations

import io
import logging
from typing import Any, Iterator, Optional

from docx import Document as DocxDocument
from docx.document import Document as DocxDocumentType
from docx.table import _Cell
from pptx import Presentation

logger = logging.getLogger(__name__)


# ── Structural heuristics ───────────────────────────────────────────────────
#
# Anything flagged is_structural=True is returned as-is by the rewriter —
# Claude isn't allowed to touch these. They're the skeleton of the document:
# section headings, labels, table column headers, etc. Everything else is
# fair game for content rewriting.

# Known doc "skeleton" words that should never be rewritten — case-insensitive
# full-text match after stripping punctuation.
_STRUCTURAL_LABELS = {
    "attendees",
    "next steps",
    "action item",
    "action items",
    "owner",
    "owners",
    "timeline",
    "collateral shared",
    "collateral",
    "note",
    "notes",
    "date",
    "dates",
    "signatures",
    "signature",
    "purpose",
    "term",
    "confidential information",
    "obligations",
    "governing law",
    "witness",
    "disclosing party",
    "receiving party",
    "executive summary",
    "key discussion points",
    "discussion points",
    "decisions",
    "decisions made",
    "name:",
    "title:",
    "by:",
    # MOM template sub-headings
    "overview",
    "key challenges",
    "current tooling",
    "areas of strong interest",
}


def _normalize_label(text: str) -> str:
    """Lowercase + strip surrounding punctuation for label matching."""
    return (text or "").strip().strip(":—-–*•").strip().lower()


def _looks_structural(style_name: str, text: str) -> bool:
    """Return True if a paragraph is a heading/label and should NOT be rewritten.

    Note on empty paragraphs: in real templates, an empty paragraph directly
    after a heading like "Overview" or "Attendees" is the *content slot*, not
    a spacer. Marking it structural here would lock Claude out of filling it,
    which is exactly the bug we hit on the Riskcovry MOM. We now treat empty
    paragraphs as content; the rewriter's prompt tells Claude to leave true
    spacers blank when no input maps to them.
    """
    # Word styles — only true headings (Heading 1–6) are skeleton. Title and
    # Subtitle styles often hold client-specific text (e.g. the meeting title
    # banner "Riskcovry and Beacon" / "Meeting Recap — 21 April 2026") that
    # MUST be rewritable. Locking those down was making the title row stale.
    style = (style_name or "").lower()
    if "heading" in style:
        return True

    if not text or not text.strip():
        # Empty paragraph with no heading style → treat as a fillable slot.
        return False

    # Known structural labels (case-insensitive full match).
    normalized = _normalize_label(text)
    if normalized in _STRUCTURAL_LABELS:
        return True

    # Very short all-caps lines are usually labels (e.g. "ATTENDEES", "TERM").
    stripped = text.strip()
    words = stripped.split()
    if len(words) <= 3 and stripped == stripped.upper() and any(c.isalpha() for c in stripped):
        return True

    return False


# ── DOCX extraction & rewriting ─────────────────────────────────────────────


def _iter_docx_paragraphs(doc: DocxDocumentType) -> Iterator[tuple[str, dict, Any]]:
    """Walk every paragraph in the doc and yield (block_type, coords, para).

    ``coords`` is a dict with table/row/col/para-index when applicable, so
    the rewriter can locate the same paragraph on a second pass.
    """
    # 1. Body paragraphs
    for para in doc.paragraphs:
        yield "paragraph", {}, para

    # 2. Table cells (and nested tables, one level deep — matches the NDA walker)
    for t_idx, table in enumerate(doc.tables):
        for r_idx, row in enumerate(table.rows):
            for c_idx, cell in enumerate(row.cells):
                for p_idx, para in enumerate(cell.paragraphs):
                    yield "table_cell", {
                        "table_index": t_idx,
                        "row_index": r_idx,
                        "col_index": c_idx,
                        "para_index_in_cell": p_idx,
                    }, para

    # 3. Headers / footers
    for s_idx, section in enumerate(doc.sections):
        for kind, container in (("header", section.header), ("footer", section.footer)):
            for p_idx, para in enumerate(container.paragraphs):
                yield kind, {
                    "section_index": s_idx,
                    "para_index": p_idx,
                }, para


def extract_docx_structure(docx_bytes: bytes) -> list[dict]:
    """Return every text-bearing block in document order, tagged for rewriting.

    Each entry has:
      block_type           "paragraph" | "table_cell" | "header" | "footer"
      block_index          sequential int across the whole document
      table_index / ...    coordinates when block_type == "table_cell"
      section_index / ...  coordinates for header/footer
      style                paragraph style name
      text                 full concatenated text of all runs
      is_structural        True = skeleton/heading, leave untouched
    """
    doc = DocxDocument(io.BytesIO(docx_bytes))
    structure: list[dict] = []
    for block_index, (block_type, coords, para) in enumerate(_iter_docx_paragraphs(doc)):
        text = "".join(run.text for run in para.runs) if para.runs else para.text
        style_name = para.style.name if para.style else ""
        entry: dict[str, Any] = {
            "block_type": block_type,
            "block_index": block_index,
            "table_index": coords.get("table_index"),
            "row_index": coords.get("row_index"),
            "col_index": coords.get("col_index"),
            "para_index_in_cell": coords.get("para_index_in_cell"),
            "section_index": coords.get("section_index"),
            "para_index": coords.get("para_index"),
            "style": style_name,
            "text": text,
            "is_structural": _looks_structural(style_name, text),
        }
        structure.append(entry)
    return structure


def _set_paragraph_text(para: Any, new_text: str) -> None:
    """Overwrite a paragraph's text while preserving the first run's formatting.

    python-docx paragraphs store text as a sequence of runs. Setting
    ``paragraph.text = "..."`` wipes all run-level formatting. Instead we
    keep the first run (inheriting its font/color/bold) and blank every
    other run.
    """
    if not para.runs:
        # Paragraph has no runs — safe to assign directly.
        para.text = new_text or ""
        return
    para.runs[0].text = new_text or ""
    for run in para.runs[1:]:
        run.text = ""


def rewrite_docx_content(
    docx_bytes: bytes,
    rewritten_blocks: list[dict],
) -> bytes:
    """Apply rewritten text to a .docx without touching any formatting.

    ``rewritten_blocks`` is a list of dicts, each with ``block_index`` and
    ``new_text``. Blocks whose index isn't in the list are left untouched.
    """
    # Build a lookup so we can skip blocks Claude didn't rewrite.
    patches: dict[int, str] = {}
    for block in rewritten_blocks or []:
        if "block_index" not in block or "new_text" not in block:
            continue
        try:
            patches[int(block["block_index"])] = str(block["new_text"] or "")
        except (TypeError, ValueError):
            logger.warning("Skipping malformed block: %r", block)

    doc = DocxDocument(io.BytesIO(docx_bytes))
    for block_index, (_block_type, _coords, para) in enumerate(_iter_docx_paragraphs(doc)):
        if block_index not in patches:
            continue
        new_text = patches[block_index]
        # Defensive guard: if Claude returned an empty string for a block
        # that originally had real text, that's almost always a mistake (a
        # hallucinated "wipe this row" rather than a deliberate clear). We
        # only allow empty rewrites when the slot was already empty.
        original_text = "".join(r.text for r in para.runs).strip()
        if new_text.strip() or not original_text:
            _set_paragraph_text(para, new_text)
        # else: leave original text untouched

    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


# ── PPTX extraction & rewriting ─────────────────────────────────────────────


def _iter_pptx_paragraphs(prs: Presentation) -> Iterator[tuple[dict, Any]]:
    """Walk every paragraph inside every text frame on every slide."""
    for s_idx, slide in enumerate(prs.slides):
        for sh_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            for p_idx, para in enumerate(tf.paragraphs):
                yield {
                    "slide_index": s_idx,
                    "shape_index": sh_idx,
                    "para_index": p_idx,
                }, para


def extract_pptx_structure(pptx_bytes: bytes) -> list[dict]:
    """Return every paragraph across every slide in order."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    structure: list[dict] = []
    for block_index, (coords, para) in enumerate(_iter_pptx_paragraphs(prs)):
        text = "".join(run.text for run in para.runs) if para.runs else ""
        entry: dict[str, Any] = {
            "block_type": "slide_text",
            "block_index": block_index,
            "slide_index": coords["slide_index"],
            "shape_index": coords["shape_index"],
            "para_index": coords["para_index"],
            "text": text,
            "is_structural": _looks_structural("", text),
        }
        structure.append(entry)
    return structure


def _set_pptx_paragraph_text(para: Any, new_text: str) -> None:
    """Same pattern as the docx helper — keep first run, blank the rest."""
    if not para.runs:
        # Empty paragraph — add one run carrying the new text.
        run = para.add_run()
        run.text = new_text or ""
        return
    para.runs[0].text = new_text or ""
    for run in para.runs[1:]:
        run.text = ""


def rewrite_pptx_content(
    pptx_bytes: bytes,
    rewritten_blocks: list[dict],
) -> bytes:
    """Apply rewritten text to a .pptx, preserving fonts/colors/layouts."""
    patches: dict[int, str] = {}
    for block in rewritten_blocks or []:
        if "block_index" not in block or "new_text" not in block:
            continue
        try:
            patches[int(block["block_index"])] = str(block["new_text"] or "")
        except (TypeError, ValueError):
            logger.warning("Skipping malformed pptx block: %r", block)

    prs = Presentation(io.BytesIO(pptx_bytes))
    for block_index, (_coords, para) in enumerate(_iter_pptx_paragraphs(prs)):
        if block_index not in patches:
            continue
        _set_pptx_paragraph_text(para, patches[block_index])

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
